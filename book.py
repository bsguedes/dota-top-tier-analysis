from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import calendar


class Book:
    def __init__(self, team_name, years, tit, sub, players, month=None):
        self.presentation = Presentation()
        self.page_number = 0
        self.team_name = team_name
        self.players = players
        self.presentation.slide_width = Inches(8.5)
        self.presentation.slide_height = Inches(11)
        y = '_'.join([str(x) for x in years])
        self.dates = y if month is None else "%s_%s" % (y, calendar.month_abbr[month])
        self.add_divider_page(tit, "%s" % sub)

    def add_divider_page(self, text, sub_text):
        slide = self.add_page(0, Color.LIGHT_GREY, show_page_number=False)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = text
        subtitle.text = sub_text

    def add_page(self, layout, color, show_page_number=True):
        slide_layout = self.presentation.slide_layouts[layout]
        slide = self.presentation.slides.add_slide(slide_layout)
        if layout == 5:
            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
        self.page_number += 1
        # if show_page_number:
        #    Slides.text_box(slide, str(self.page_number), 0.1, 0.1, font_size=10)
        return slide

    def save(self):
        self.presentation.save('report_book.pptx')


class Color:
    LIGHT_GREY = [230, 230, 230]
