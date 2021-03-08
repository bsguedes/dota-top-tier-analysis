import math
import os
import os.path
from pptx import Presentation
from pptx.util import Inches, Pt
from tier import TierItem
from category import Category
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from constants import roles
from constants import sequence
from constants import win_rate
import calendar


class Slides:
    def __init__(self, team_name, years, tit, sub, players, month=None):
        self.presentation = Presentation()
        self.page_number = 0
        self.team_name = team_name
        self.players = players
        y = '_'.join([str(x) for x in years])
        self.dates = y if month is None else "%s_%s" % (y, calendar.month_abbr[month])
        self.add_divider_slide(tit, "%s" % sub)

    def add_slide(self, layout, r, g, b, show_page_number=True):
        slide_layout = self.presentation.slide_layouts[layout]
        slide = self.presentation.slides.add_slide(slide_layout)
        if layout == 5:
            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        Slides.change_slide_color(slide, r, g, b)
        self.page_number += 1
        if show_page_number:
            Slides.text_box(slide, str(self.page_number), 0.1, 0.1,
                            font_size=10)
        return slide

    def add_divider_slide(self, text, sub_text):
        slide = self.add_slide(0, 230, 230, 230, show_page_number=False)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = text
        subtitle.text = sub_text

    def add_match_roles_summary(self, data):
        top = 1.4
        left = 0.3
        pic_size = 0.25
        row = 0.28
        rows = 20
        spacing = 0.7

        for c in range(math.ceil(len(data) / rows)):
            slide = self.add_slide(5, 244, 244, 244)
            title_shape = slide.shapes.title
            title_shape.text = 'Match Summary with Roles'

            Slides.text_box(slide, 'Match ID', left, top, font_size=12, bold=True)
            Slides.text_box(slide, 'Win', left + 1.1, top, font_size=12, bold=True)
            Slides.text_box(slide, 'Duration', left + 1.5, top, font_size=12, bold=True)
            Slides.text_box(slide, 'Hard Carry', left + 2.4, top, font_size=12)
            Slides.text_box(slide, 'Mid', left + 2.4 + 1 * (spacing + pic_size), top, font_size=12)
            Slides.text_box(slide, 'Offlane', left + 2.4 + 2 * (spacing + pic_size), top, font_size=12)
            Slides.text_box(slide, 'Support', left + 2.4 + 3 * (spacing + pic_size), top, font_size=12)
            Slides.text_box(slide, 'Hard Support', left + 2.4 + 4 * (spacing + pic_size), top, font_size=12)
            Slides.text_box(slide, 'Lanes', left + 2.4 + 5 * (spacing + pic_size) + 0.5, top, font_size=12, bold=True)

            for r in range(rows):
                index = rows * c + r
                if len(data) > index:
                    match = data[index]
                    txt_box = slide.shapes.add_textbox(Inches(left), Inches(top + row * (r + 1)),
                                                       Inches(row), Inches(row))
                    tf = txt_box.text_frame
                    Slides.hyperlink_sequence(tf.paragraphs[0], [str(int(match['match_id']))], 1)
                    tf.paragraphs[0].font.size = Pt(12)
                    Slides.text_box(slide, 'yes' if match['win'] else 'no', left + 1.1, top + row * (r + 1),
                                    font_size=11)
                    Slides.text_box(slide, '%02d:%02d' % (match['duration'] // 60, match['duration'] % 60),
                                    left + 1.6, top + row * (r + 1), font_size=11)
                    for player in match['players']:
                        for pos, i in zip(['hard carry', 'mid', 'offlane', 'support', 'hard support'], range(5)):
                            if player['position'] == pos:
                                pic_path = 'data/pics/%s.jpg' % player['player_id']
                                if os.path.isfile(pic_path):
                                    slide.shapes.add_picture(pic_path, Inches(left + 2.5 + i * (spacing + pic_size)),
                                                             Inches(top + row * (r + 1)), height=Inches(pic_size))
                                hero_path = 'data/heroes/%s.jpg' % player['hero_id']
                                if os.path.isfile(hero_path):
                                    slide.shapes.add_picture(hero_path,
                                                             Inches(left + 2.5 + i * (spacing + pic_size) + row),
                                                             Inches(top + row * (r + 1)), height=Inches(pic_size))
                    Slides.text_box(slide, str(match['lane']), left + 2.4 + 5 * (spacing + pic_size) + 0.5,
                                    top + row * (r + 1), font_size=11)

    def add_item_slides(self, item_data):
        for data in item_data:
            slide = self.add_slide(5, 214, 187, 242)
            title_shape = slide.shapes.title
            title_shape.text = data['item_name']
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            pic_path = 'data/items/%s.jpg' % data['item_internal_name']
            slide.shapes.add_picture(pic_path, Inches(7.5), Inches(0.2), height=Inches(1.1))
            headers = ['Amount', 'Wins', 'Losses', 'Matches', 'Win Rate']
            keys = ['amount', 'wins', 'losses', 'matches', 'wr']
            formats = ['%s', '%s', '%s', '%s', '%.2f %%']
            table_data = [{'amount': a, 'wins': b, 'losses': c, 'matches': d, 'wr': e} for a, b, c, d, e in
                          zip(data['counts'], data['wins'], data['losses'], data['matches'], data['wr'])]
            Slides.create_table(slide, table_data, headers, keys, formats, Inches(2), Inches(1.8), Inches(6),
                                1, 12, 15)
            chart_data = ChartData()
            chart_data.categories = data['counts']
            chart_data.add_series('Wins', data['wins'])
            chart_data.add_series('Losses', data['losses'])
            graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(4.2), Inches(9),
                                                   Inches(3.2),
                                                   chart_data)
            chart = graphic_frame.chart
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(10)
            data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
            data_labels.position = XL_LABEL_POSITION.INSIDE_END
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False

    def add_heroes(self, hero_stats, min_matches_with_hero):
        for hero in hero_stats:
            if hero['matches'] > 0:
                slide = self.add_slide(5, 204, 255, 255)
                title_shape = slide.shapes.title
                title_shape.text = '%s' % hero['name']
                title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                pic_path = 'data/heroes/%s.jpg' % hero['id']
                slide.shapes.add_picture(pic_path, Inches(7), Inches(0.2), height=Inches(1.1))

                headers = ['Role', 'Matches', 'Win Rate']
                keys = ['role', 'matches', 'wr']
                formats = ['%s', '%s', '%.2f %%']
                widths = [1.5, 1, 1]
                Slides.create_table(slide, hero['roles'], headers, keys, formats, Inches(0.5), Inches(1.5), Inches(3.5),
                                    1,
                                    12, 15, widths=widths)

                tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(1.5), Inches(0.5))
                tf = tx_box.text_frame
                p = tf.paragraphs[0]
                p.text = '%s has %s matches' % (hero['name'], hero['matches'])

                s = sorted([x for x in hero['played_by'] if x['matches'] >= min_matches_with_hero and x['rating'] > 0],
                           key=lambda e: (e['rating'], e['matches']),
                           reverse=True)
                if len(s) > 0:
                    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(1.5), Inches(0.5))
                    tf = tx_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = 'Best players with at least %s matches:' % min_matches_with_hero
                    p.font.size = Pt(16)
                    for i in range(min(3, len(s))):
                        pic_path = 'data/pics/%s.jpg' % self.players[s[i]['name']]
                        if os.path.isfile(pic_path):
                            slide.shapes.add_picture(pic_path, Inches(0.25), Inches(4.8) + i * Inches(0.8),
                                                     height=Inches(0.7))

                        tx_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.85) + i * Inches(0.8), Inches(1.5),
                                                          Inches(0.5))
                        p = tx_box.text_frame.paragraphs[0]
                        p.text = '%.2f' % s[i]['rating']
                        p.font.size = Pt(24)
                        p.font.bold = True

                        tx_box = slide.shapes.add_textbox(Inches(2), Inches(5.05) + i * Inches(0.8), Inches(1.5),
                                                          Inches(0.5))
                        p = tx_box.text_frame.paragraphs[0]
                        p.text = '%i-%i (%.2f %% wr)' % (s[i]['wins'], s[i]['matches'] - s[i]['wins'], s[i]['wr'])
                        p.font.size = Pt(15)

                        tx_box = slide.shapes.add_textbox(Inches(2), Inches(4.75) + i * Inches(0.8), Inches(1.5),
                                                          Inches(0.5))
                        p = tx_box.text_frame.paragraphs[0]
                        p.text = '%s' % s[i]['name']
                        p.font.size = Pt(14)
                        p.font.bold = True

                headers = ['Played by', 'Matches', 'Wins', 'Win Rate', 'Rating']
                keys = ['name', 'matches', 'wins', 'wr', 'rating']
                formats = ['%s', '%s', '%s', '%.2f %%', '%.2f']
                widths = [1.3, 1, 0.8, 1, 0.9]
                Slides.create_table(slide, hero['played_by'], headers, keys, formats, Inches(4.5), Inches(1.5),
                                    Inches(5), 1, 11, 15, widths=widths)

    def add_player_activity_data(self, desc):
        slide = self.add_slide(5, 255, 229, 204)
        title_shape = slide.shapes.title
        title_shape.text = desc['name']
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        pic_path = 'data/pics/%s.jpg' % desc['id']
        if os.path.isfile(pic_path):
            slide.shapes.add_picture(pic_path, Inches(8), Inches(0.2), height=Inches(1.1))

        data = desc['months']
        chart_data = ChartData()
        chart_data.categories = data.keys()
        chart_data.add_series('Wins', [o['wins'] for o in data.values()])
        chart_data.add_series('Losses', [o['losses'] for o in data.values()])
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(9),
                                               Inches(5.5),
                                               chart_data)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    def add_player_tables_slide(self, desc, fantasy=None):
        slide = self.add_slide(5, 255, 229, 204)
        title_shape = slide.shapes.title
        title_shape.text = desc['name']
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        pic_path = 'data/pics/%s.jpg' % desc['id']
        if os.path.isfile(pic_path):
            slide.shapes.add_picture(pic_path, Inches(8), Inches(0.2), height=Inches(1.1))

        headers = ['Role', 'Matches', 'Win Rate']
        keys = ['role', 'matches', 'wr']
        formats = ['%s', '%s', '%.2f %%']
        widths = [1.5, 1, 1]

        if len(fantasy) > 0:
            headers = ['Role', 'Matches', 'Win Rate', 'PnKoins']
            keys = ['role', 'matches', 'wr', 'coins']
            formats = ['%s', '%s', '%.2f %%', '%s']
            widths = [1.1, 1, 1, 1]
            for role in desc['roles']:
                coins = 0 if role['role'] not in fantasy else fantasy[role['role']]
                role['coins'] = coins if coins > 0 else ' '

        Slides.create_table(slide, desc['roles'], headers, keys, formats, Inches(0.5), Inches(1.5), Inches(3.5), 1, 12,
                            15,
                            widths=widths)

        headers = ['Paired with', 'Matches', 'Wins', 'Win Rate']
        keys = ['name', 'matches', 'wins', 'wr']
        formats = ['%s', '%s', '%s', '%.2f %%']
        widths = [1.5, 1, 1, 1]
        Slides.create_table_with_text_boxes(slide, desc['pairings'], headers, keys, formats, 5, 1.5, 4.5,
                                            11, 15, widths=widths, line_spacing=0.2)

        heroes = desc['top_heroes']
        if len(heroes) > 0:
            Slides.text_box(slide, 'Best %s player with (based on hero rating):' % self.team_name, 0.3, 3.4,
                            font_size=13)
            rows = 7
            columns = 4
            for y in range(columns):
                for x in range(rows):
                    index = y * rows + x
                    if index < len(heroes):
                        hero = heroes[index]
                        pic_path = 'data/heroes/%s.jpg' % hero
                        slide.shapes.add_picture(pic_path, Inches(0.3) + y * Inches(3.8 / columns),
                                                 Inches(3.8) + x * Inches(0.5),
                                                 height=Inches(0.45))

    def add_player_summary(self, player_descriptors, min_matches):
        slide = self.add_slide(5, 222, 200, 178)
        title_shape = slide.shapes.title
        title_shape.text = 'Player Summary'

        win_rate_lambda = lambda d: win_rate(d['wins'], d['matches'])
        win_rate_without_lambda = lambda d: win_rate(d['team_wins'] - d['wins'], d['team_matches'] - d['matches'])
        losses_lambda = lambda d: d['matches'] - d['wins']
        hero_count_lambda = lambda d: len([1 for _, h in d['heroes'].items() if h['matches'] > 0])

        headers = ['Name', 'Rating', 'Matches', 'Win Rate', 'Win Rate w/o',
                   'Wins', 'Losses', 'Heroes', 'Versat.', 'MMR Change']
        keys = ['name', 'rating', 'matches', win_rate_lambda, win_rate_without_lambda,
                'wins', losses_lambda, hero_count_lambda, 'versatility', 'mmr_var']
        formats = ['%s', '%.2f', '%s', '%.2f %%', '%.2f %%', '%s', '%s', '%s', '%.3f', '%s']
        widths = [1.4, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9]

        players = sorted([p for p in player_descriptors if p['matches'] >= min_matches],
                         key=lambda e: e['rating'],
                         reverse=True)
        Slides.create_table(slide, players, headers, keys, formats, Inches(0.25), Inches(1.5), Inches(9.5), 1, 12, 14,
                            widths=widths)

    def add_player_mmr(self, mmr_changes):
        slide = self.add_slide(5, 222, 200, 178)
        title_shape = slide.shapes.title
        title_shape.text = 'MMR diff in ranked %s matches' % self.team_name

        if len(mmr_changes) > 0:
            chart_data = CategoryChartData()
            chart_data.categories = [p['name'] for p in mmr_changes]
            chart_data.add_series('MMR', [p['mmr'] for p in mmr_changes])

            x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(6)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(12)
            category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = Pt(12)

            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(12)
            data_labels.font.color.rgb = RGBColor(0, 0, 0)

    def add_player_data_slide(self, desc, min_matches_with_hero):
        slide = self.add_slide(5, 255, 229, 204)
        title_shape = slide.shapes.title
        title_shape.text = desc['name']
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        pic_path = 'data/pics/%s.jpg' % desc['id']
        if os.path.isfile(pic_path):
            slide.shapes.add_picture(pic_path, Inches(8), Inches(0.2), height=Inches(1.1))

        Slides.text_box(slide, 'Rating:', 0.5, 1.6, 1.5)
        Slides.text_box(slide, 'Wins:', 0.5, 2, 1.5)
        Slides.text_box(slide, 'Losses:', 0.5, 2.4, 1.5)
        Slides.text_box(slide, 'Distinct heroes:', 0.5, 2.8, 1.5)
        Slides.text_box(slide, 'Matches played:', 0.5, 3.2, 1.5)

        Slides.text_box(slide, str(desc['wins']), 1.6, 1.95, 1.75, alignment=PP_ALIGN.RIGHT, font_size=22)
        Slides.text_box(slide, str(desc['matches'] - desc['wins']), 1.6, 2.35, 1.75, alignment=PP_ALIGN.RIGHT,
                        font_size=22)
        Slides.text_box(slide, str(len([x for x in desc['heroes'].values() if x['matches'] > 0])), 1.6, 2.75, 1.75,
                        alignment=PP_ALIGN.RIGHT, font_size=22)
        Slides.text_box(slide, str(desc['matches']), 1.6, 3.15, 1.75, alignment=PP_ALIGN.RIGHT, font_size=22)

        Slides.text_box(slide, 'Best Win Streak:', 3.6, 1.6, 1.5)
        Slides.text_box(slide, 'Worst Loss Streak:', 3.6, 2, 1.5)
        Slides.text_box(slide, 'Radiant Wins:', 3.6, 2.8, 1.5)
        Slides.text_box(slide, 'Dire Wins:', 3.6, 3.2, 1.5)

        Slides.text_box(slide, str(max(0, max(desc['streaks']))), 4.8, 1.55, 1.75, alignment=PP_ALIGN.RIGHT,
                        font_size=22)
        Slides.text_box(slide, str(abs(min(0, min(desc['streaks'])))), 4.8, 1.95, 1.75, alignment=PP_ALIGN.RIGHT,
                        font_size=22)
        Slides.text_box(slide, '%.2f %%' % desc['radiant_wr'], 4.8, 2.75, 1.75, alignment=PP_ALIGN.RIGHT, font_size=22)
        Slides.text_box(slide, '%.2f %%' % desc['dire_wr'], 4.8, 3.15, 1.75, alignment=PP_ALIGN.RIGHT, font_size=22)

        Slides.text_box(slide, 'Win Rate:', 6.7, 1.6, 1.5)
        Slides.text_box(slide, 'Versatility:', 6.7, 2, 1.5)

        tx_box = slide.shapes.add_textbox(Inches(1.6), Inches(1.5), Inches(1.75), Inches(0.5))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = '%.2f' % desc['rating']
        p.font.bold = True
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.RIGHT

        tx_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.55), Inches(1.75), Inches(0.5))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = '%.2f %%' % (100 * desc['wins'] / desc['matches'])
        p.font.size = Pt(22)
        p.alignment = PP_ALIGN.RIGHT

        tx_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.95), Inches(1.75), Inches(0.5))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = '%.3f' % desc['versatility']
        p.font.size = Pt(22)
        p.alignment = PP_ALIGN.RIGHT

        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(1.5), Inches(0.5))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = 'Best heroes by rating'
        p.font.size = Pt(16)

        filtered_heroes = [e for e in desc['heroes'].items() if e[1]['matches'] >= min_matches_with_hero]

        heroes = sorted(filtered_heroes, key=lambda e: e[1]['rating'], reverse=True)
        for i in range(4):
            if len(heroes) > i and heroes[i][1]['rating'] > 0:
                pic_path = 'data/heroes/%s.jpg' % heroes[i][0]
                slide.shapes.add_picture(pic_path, Inches(0.5), Inches(4.2) + i * Inches(0.8), height=Inches(0.7))
                tx_box = slide.shapes.add_textbox(Inches(1.9), Inches(4.3) + i * Inches(0.8), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = '%.2f' % heroes[i][1]['rating']
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(24)
                tx_box = slide.shapes.add_textbox(Inches(2.65), Inches(4.2) + i * Inches(0.8), Inches(1.5), Inches(0.5))
                tf = tx_box.text_frame
                p = tf.paragraphs[0]
                p.text = '%s matches' % heroes[i][1]['matches']
                p.font.size = Pt(16)
                tx_box = slide.shapes.add_textbox(Inches(2.65), Inches(4.45) + i * Inches(0.8), Inches(1), Inches(0.5))
                tf = tx_box.text_frame
                p = tf.paragraphs[0]
                p.text = '%s - %s' % (heroes[i][1]['wins'], heroes[i][1]['matches'] - heroes[i][1]['wins'])
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.CENTER

        tx_box = slide.shapes.add_textbox(Inches(4), Inches(3.7), Inches(1.5), Inches(0.5))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = 'Most played heroes'
        p.font.size = Pt(16)
        s = sorted(desc['heroes'].items(), key=lambda e: e[1]['matches'], reverse=True)
        for i in range(4):
            if len(s) > i and s[i][1]['matches'] > 0:
                pic_path = 'data/heroes/%s.jpg' % s[i][0]
                slide.shapes.add_picture(pic_path, Inches(4), Inches(4.2) + i * Inches(0.8), height=Inches(0.7))
                tx_box = slide.shapes.add_textbox(Inches(5.5), Inches(4.3) + i * Inches(0.8), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = '%i' % s[i][1]['matches']
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(24)

        tx_box = slide.shapes.add_textbox(Inches(6.5), Inches(3.7), Inches(1.5), Inches(0.5))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = 'Worst rating against'
        p.font.size = Pt(16)
        s = sorted(desc['heroes'].items(), key=lambda e: e[1]['inv_rating'], reverse=True)
        for i in range(4):
            if len(s) > i and s[i][1]['matches_against'] > 0:
                pic_path = 'data/heroes/%s.jpg' % s[i][0]
                slide.shapes.add_picture(pic_path, Inches(6.5), Inches(4.2) + i * Inches(0.8), height=Inches(0.7))
                tx_box = slide.shapes.add_textbox(Inches(7.9), Inches(4.3) + i * Inches(0.8), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = '%.2f' % s[i][1]['inv_rating']
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(24)
                tx_box = slide.shapes.add_textbox(Inches(8.65), Inches(4.2) + i * Inches(0.8), Inches(1.5), Inches(0.5))
                tf = tx_box.text_frame
                p = tf.paragraphs[0]
                p.text = '%s matches' % s[i][1]['matches_against']
                p.font.size = Pt(16)
                tx_box = slide.shapes.add_textbox(Inches(8.65), Inches(4.45) + i * Inches(0.8), Inches(1), Inches(0.5))
                tf = tx_box.text_frame
                p = tf.paragraphs[0]
                p.text = '%s - %s' % (s[i][1]['wins_against'], s[i][1]['matches_against'] - s[i][1]['wins_against'])
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.CENTER

    def add_intro_slide(self, match_count, min_player_count, min_matches, min_couples_matches):
        slide = self.add_slide(1, 152, 251, 152)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = '%s Summary' % self.team_name
        tf = body_shape.text_frame
        tf.text = 'A total of %i matches were analyzed.' % match_count
        p = tf.add_paragraph()
        p.text = 'Parties with at least %i %s players.' % (min_player_count, self.team_name)
        p.level = 1
        p = tf.add_paragraph()
        p.text = 'Average metrics consider players with at least %i matches.' % min_matches
        p.level = 1
        p = tf.add_paragraph()
        p.text = 'For best pairings and best players in roles, a minimum of %i matches.' % min_couples_matches
        p.level = 2
        p.font.size = Pt(20)
        p = tf.add_paragraph()
        p.text = 'Best team per hero on each role requires a minimum of %i matches for that hero.' % min_couples_matches
        p.level = 2
        p.font.size = Pt(20)
        p = tf.add_paragraph()
        p.text = 'Best team per hero and player on each role requires a minimum of %i matches.' \
                 % (min_couples_matches / 2)
        p.level = 2
        p.font.size = Pt(20)
        p = tf.add_paragraph()
        p.text = 'The list of players is: %s' % (', '.join(self.players.keys()))
        p.level = 1
        p.font.size = Pt(18)

    def add_five_player_compositions(self, comp, matches):
        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Full Parties with >= %i matches' % matches
        headers = ['Players', 'Wins', 'Matches', 'Win Rate']
        keys = ['players', 'wins', 'matches', 'wr']
        formats = ['%s', '%s', '%s', '%.2f %%']
        widths = [4.5, 1.5, 1.5, 1.5]
        Slides.create_table(slide, comp, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths)

    def add_match_details(self, to_parse, match_types, skills):
        slide = self.add_slide(1, 152, 251, 152)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = '%s Summary' % self.team_name
        tf = body_shape.text_frame
        tf.text = '%s matches have incomplete data' % to_parse
        p = tf.add_paragraph()
        p.font.size = Pt(20)
        p.text = 'Login every week to OpenDota in order to have complete data to our matches'
        p = tf.add_paragraph()
        p.text = 'Match Summary per Game Type:'
        headers = ['Game Type', 'Wins', 'Matches', 'Win Rate']
        keys = ['lobby_type', 'wins', 'matches', 'wr']
        formats = ['%s', '%s', '%s', '%.2f %%']
        widths = [4.5, 1.5, 1.5, 1.5]
        Slides.create_table(slide, match_types, headers, keys, formats, Inches(0.5), Inches(4), Inches(9), 1, 13, 15,
                            widths=widths)

        slide = self.add_slide(1, 152, 251, 152)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = '%s Summary' % self.team_name
        tf = body_shape.text_frame
        tf.text = 'Match Summary per Game Skill'
        headers = ['Game Skill', 'Wins', 'Matches', 'Win Rate']
        keys = ['skill', 'wins', 'matches', 'wr']
        formats = ['%s', '%s', '%s', '%.2f %%']
        widths = [4.5, 1.5, 1.5, 1.5]
        Slides.create_table(slide, skills, headers, keys, formats, Inches(0.5), Inches(2.5), Inches(9), 1, 13, 15,
                            widths=widths)

    def add_win_rate_details_slide(self, fb_object, bounties):
        slide = self.add_slide(1, 152, 251, 152)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = '%s Summary' % self.team_name
        textbox = slide.shapes[1]
        sp = textbox.element
        sp.getparent().remove(sp)

        Slides.text_box(slide, 'First Blood rate', 0.8, 2, font_size=24, bold=True)
        Slides.text_box(slide, '%.2f %%' % fb_object['first_blood_rate'], 4.5, 2, font_size=24)
        Slides.text_box(slide, '%s matches' % fb_object['matches'], 6, 2, font_size=24)
        Slides.text_box(slide, '%s-%s' % (fb_object['first_bloods'], fb_object['matches'] - fb_object['first_bloods']),
                        8, 2, font_size=24)

        Slides.text_box(slide, 'Wins with first blood', 0.8, 2.5, font_size=24, bold=True)
        Slides.text_box(slide, '%.2f %%' % fb_object['wr_if_first_blood'], 4.5, 2.5, font_size=24)
        Slides.text_box(slide, '%s matches' % fb_object['first_bloods'], 6, 2.5, font_size=24)
        Slides.text_box(slide, '%s-%s' % (
                                            fb_object['wins_if_first_blood'],
                                            fb_object['first_bloods'] - fb_object['wins_if_first_blood']),
                        8, 2.5, font_size=24)

        Slides.text_box(slide, 'Wins without first blood', 0.8, 3, font_size=24, bold=True)
        Slides.text_box(slide, '%.2f %%' % fb_object['wr_if_not_first_blood'], 4.5, 3, font_size=24)
        Slides.text_box(slide, '%s matches' % (fb_object['matches'] - fb_object['first_bloods']), 6, 3,
                        font_size=24)
        Slides.text_box(slide, '%s-%s' % (fb_object['wins_if_no_first_blood'],
                                          fb_object['matches'] - fb_object['first_bloods'] - fb_object[
                                              'wins_if_no_first_blood']),
                        8, 3, font_size=24)

        Slides.text_box(slide, 'Bounty runes at minute zero', 0.8, 3.8, font_size=24, bold=True)
        headers = ['Bounties', 'Wins', 'Losses', 'Matches', 'Win Rate']
        keys = ['counts', 'wins', 'losses', 'matches', 'wr']
        formats = ['%s', '%s', '%s', '%s', '%.2f %%']
        Slides.create_table(slide, bounties, headers, keys, formats, Inches(2), Inches(4.5), Inches(6),
                            1, 12, 15)

    def add_win_rate_slide(self, winrate, match_count, party_size, faction):
        slide = self.add_slide(1, 152, 251, 152)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = '%s Summary' % self.team_name
        tf = body_shape.text_frame
        tf.text = 'Win rate: %.2f %%' % winrate
        p = tf.add_paragraph()
        p.text = 'Match count: %s' % match_count
        tf.add_paragraph()
        for i in range(1, 6):
            p = tf.add_paragraph()
            p.font.size = Pt(20)
            p.text = '%s matches with party size = %i' % (party_size[i-1], i)
        Slides.text_box(slide, 'Radiant', 0.8, 6, font_size=24, bold=True)
        Slides.text_box(slide, 'Dire', 0.8, 6.6, font_size=24, bold=True)
        Slides.text_box(slide, '%.2f %%' % faction['r_wr'], 2.5, 6, font_size=24)
        Slides.text_box(slide, '%.2f %%' % faction['d_wr'], 2.5, 6.6, font_size=24)
        Slides.text_box(slide, '%s matches' % (faction['r_win'] + faction['r_loss']), 4, 6, font_size=24)
        Slides.text_box(slide, '%s matches' % (faction['d_win'] + faction['d_loss']), 4, 6.6, font_size=24)
        Slides.text_box(slide, '%s-%s' % (faction['r_win'], faction['r_loss']), 6, 6, font_size=24)
        Slides.text_box(slide, '%s-%s' % (faction['d_win'], faction['d_loss']), 6, 6.6, font_size=24)

    def add_top_fifteen(self, comebacks, throws, fast_wins, fast_losses, longest):
        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Top 15 %s Comebacks' % self.team_name

        headers = ['Match ID', 'Gold', 'Players']
        keys = ['match', 'gold', 'players']
        formats = ['%i', '%s', '%s']
        widths = [2, 2, 5]
        Slides.create_table(slide, comebacks, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths, hyperlink=[0])

        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Top 15 %s Throws' % self.team_name
        Slides.create_table(slide, throws, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths, hyperlink=[0])

        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        headers = ['Match ID', 'Time [min]', 'Players']
        keys = ['match', 'time', 'players']
        formats = ['%i', '%02d:%02d', '%s']
        title_shape.text = 'Top 15 %s Fast Wins (no abandons)' % self.team_name
        Slides.create_table(slide, fast_wins, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths, hyperlink=[0])

        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Top 15 %s Fast Losses (no abandons)' % self.team_name
        Slides.create_table(slide, fast_losses, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths, hyperlink=[0])

        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Top 15 %s Longest Matches' % self.team_name
        headers = ['Match ID', 'Gold', 'Players', 'Win?']
        keys = ['match', 'time', 'players', 'win']
        formats = ['%i', '%02d:%02d', '%s', '%s']
        widths = [1.7, 1.7, 3.9, 1.7]
        Slides.create_table(slide, longest, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths, hyperlink=[0])

    def add_advantage_chart(self, data, type):
        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = '%s average minute performance [%%]' % self.team_name

        data = {a: int(b * 10000) / 100 for a, b in zip(range(len(data)), data)}

        chart_data = ChartData()
        chart_data.categories = data.keys()
        chart_data.add_series(type, data.values())
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.5), Inches(1.5), Inches(9),
                                               Inches(6),
                                               chart_data)
        chart = graphic_frame.chart
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(8)

    def add_win_rate_by_match_len(self, input_data):
        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = '%s Win Rate by Game Duration [min]' % self.team_name

        Slides.histogram_win_loss(slide, input_data, 'duration', 0.5, 1.5, 9, 5)

    def add_win_rate_by_date(self, input_data, label):
        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = '%s Win Rate by %s' % (self.team_name, label)

        data = {a: b for a, b in input_data.items() if int(a) < 12} if label == 'Hour' else input_data
        headers = [a for a in data.keys()]
        keys = [str(i) for i in range(len(data))]
        formats = ['%.1f %%'] * len(data)
        table_data = [{str(i): o['wr'] for i, o in zip(range(len(data)), data.values())}]
        Slides.create_table(slide, table_data, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 12, 15)

        if label == 'Hour':
            data = {a: b for a, b in input_data.items() if int(a) >= 12}
            headers = [a for a in data.keys()]
            keys = [str(i) for i in range(len(data))]
            formats = ['%.1f %%'] * len(data)
            table_data = [{str(i): o['wr'] for i, o in zip(range(len(data)), data.values())}]
            Slides.create_table(slide, table_data, headers, keys, formats, Inches(0.5), Inches(2.5), Inches(9), 1, 12,
                                15)

        data = input_data
        chart_data = ChartData()
        chart_data.categories = data.keys()
        chart_data.add_series('Wins', [o['wins'] for o in data.values()])
        chart_data.add_series('Losses', [o['losses'] for o in data.values()])
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(3.5), Inches(9),
                                               Inches(4),
                                               chart_data)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    def add_compositions(self, compositions):
        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Lane Compositions'
        headers = ['Lane Composition', 'Wins', 'Matches', 'Win Rate']
        keys = ['comp', 'wins', 'matches', 'wr']
        formats = ['%s', '%s', '%s', '%.2f %%']
        Slides.create_table_with_text_boxes(slide, compositions, headers, keys, formats, 0.5, 1.5, 9, 12, 14)

    def add_win_rate_heroes(self, versus, text):
        versus = [v for v in versus if v['matches'] > 0]
        for c in range(math.ceil(len(versus)/14)):
            slide = self.add_slide(5, 204, 255, 204)
            title_shape = slide.shapes.title
            title_shape.text = '%s Performance %s Heroes' % (self.team_name, text)
            for y in range(2):
                for x in range(7):
                    index = c*14 + y*7 + x
                    if len(versus) > index:
                        hero = versus[index]
                        pic_path = 'data/heroes/%s.jpg' % hero['id']
                        slide.shapes.add_picture(pic_path, Inches(0.5) + y * Inches(5), Inches(1.8) + x * Inches(0.8),
                                                 height=Inches(0.7))
                        tx_box = slide.shapes.add_textbox(Inches(0.5) + y * Inches(5) + Inches(1.5),
                                                          Inches(1.9) + x * Inches(0.8), Inches(1), Inches(0.4))
                        tf = tx_box.text_frame
                        tf.text = '%.2f' % hero['rating']
                        tf.paragraphs[0].font.bold = True
                        tf.paragraphs[0].font.size = Pt(24)

                        tx_box = slide.shapes.add_textbox(Inches(0.5) + y * Inches(5) + Inches(2.3),
                                                          Inches(2) + x * Inches(0.8), Inches(1.5),
                                                          Inches(0.5))
                        tf = tx_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = '%i-%i (%.2f %% wr)' % (hero['wins'], hero['matches'] - hero['wins'], hero['wr'])
                        p.font.size = Pt(16)

                        tx_box = slide.shapes.add_textbox(Inches(0.5) + y * Inches(5) + Inches(2.3),
                                                          Inches(1.8) + x * Inches(0.8), Inches(1.5),
                                                          Inches(0.5))
                        tf = tx_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = hero['name']
                        p.font.size = Pt(12)
                        p.font.bold = True

    def add_most_played(self, heroes, most_played):
        if not most_played and len(heroes) == 0:
            return None
        rows = 7
        columns = 2 if most_played else 6
        for c in range(math.ceil(len(heroes)/(columns*rows))):
            slide = self.add_slide(5, 176, 255, 176) if most_played else self.add_slide(5, 140, 255, 140)
            title_shape = slide.shapes.title
            title_shape.text = '%s\'s %s Heroes' % (self.team_name, 'Most Played' if most_played else 'Not Played')
            for y in range(columns):
                for x in range(rows):
                    index = c*(columns * rows) + y*rows + x
                    if len(heroes) > index:
                        hero = heroes[index]
                        pic_path = 'data/heroes/%s.jpg' % hero['id']
                        slide.shapes.add_picture(pic_path, Inches(0.5) + y * Inches(9 / columns),
                                                 Inches(1.8) + x * Inches(0.8),
                                                 height=Inches(0.7))
                        if most_played:
                            tx_box = slide.shapes.add_textbox(Inches(0.5) + y * Inches(5) + Inches(2),
                                                              Inches(1.7) + x * Inches(0.8), Inches(1.5), Inches(0.5))
                            tf = tx_box.text_frame
                            p = tf.paragraphs[0]
                            p.text = '%s matches' % (hero['matches'])
                            p.font.size = Pt(22)
                            p.alignment = PP_ALIGN.CENTER
                            tx_box = slide.shapes.add_textbox(Inches(0.5) + y * Inches(5) + Inches(2),
                                                              Inches(1.65) + x * Inches(0.8) + Inches(0.5), Inches(1.5),
                                                              Inches(0.5))
                            tf = tx_box.text_frame
                            p = tf.paragraphs[0]
                            p.text = hero['name']
                            p.font.size = Pt(12)
                            p.font.bold = True
                            p.alignment = PP_ALIGN.CENTER

    def add_match_summary_by_player(self, summary, team, party_size):
        cat = Category(0, '', unit='matches')

        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Most Matches Played'
        scores = [TierItem(summary[i]['player'], summary[i]['matches'], '') for i in range(3)]
        self.add_top_three_table(scores, slide, True, cat, Inches(1), Inches(2))

        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Most Matches Played with %s' % self.team_name
        scores = [TierItem(team[i]['player'], team[i]['team_matches'], '') for i in range(3)]
        self.add_top_three_table(scores, slide, True, cat, Inches(1), Inches(2))

        slide = self.add_slide(5, 152, 251, 152)
        title_shape = slide.shapes.title
        title_shape.text = 'Match Summary by Player'
        headers = ['Player', 'Total Matches', 'with party >= 2', '%% with %s' % self.team_name,
                   'with party >= %s' % party_size]
        keys = ['player', 'matches', 'team_matches', 'perc_with_team', 'matches_party']
        formats = ['%s', '%s', '%s', '%.2f %%', '%s']
        Slides.create_table_with_text_boxes(slide,
                                            sorted([x for x in summary if x['team_matches'] > 0],
                                                   key=lambda e: e['team_matches'], reverse=True),
                                            headers, keys, formats, 0.5, 1.35, 9, 11, 14, line_spacing=0.19)

    def add_tier_slides(self, tier, category):
        texts = tier.list_to_print()

        slide = self.add_slide(6, 255, 255, 224)
        left = top = width = height = Inches(0.4)
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        tf = txt_box.text_frame
        tf.text = texts[1]
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(22)

        txt_box = slide.shapes.add_textbox(Inches(8.5), top, width, height)
        tf = txt_box.text_frame
        tf.text = "Tier points: %s" % tier.weight
        tf.paragraphs[0].font.size = Pt(14)

        scores = tier.get_top_three()
        self.add_top_three_small_table(scores, slide, tier.is_max, category, 0.3, 1)

        if len(tier.scores_array) > 0:
            chart_data = CategoryChartData()
            chart_data.categories = [x.name for x in tier.scores_array]
            chart_data.add_series(texts[1], [x.score for x in tier.scores_array])
            x, y, cx, cy = Inches(0.1), Inches(2.8), Inches(5), Inches(4.7)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(9)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = Pt(12)

        txt_box = slide.shapes.add_textbox(Inches(5.2), Inches(0.6), width, height)
        tf = txt_box.text_frame
        for i in range(2, len(texts)):
            p = tf.add_paragraph()
            p.font.size = Pt(10)
            s = texts[i].split('#')
            if len(s) == 1:
                p.text = texts[i]
                p.font.size = Pt(10)
            else:
                for j in range(len(s)):
                    r = p.add_run()
                    r.text = s[j]
                    if j == 1:
                        Slides.add_open_dota_link(r, s[j])
            if texts[i].startswith("Tier"):
                p.font.bold = True

    def add_top_three_small_table(self, scores, slide, is_max, category, left, top):
        pos = [(0, '1'), (1, '2'), (2, '3')]
        spacing = 1.6
        val = min(3, len(scores))
        for i, name in pos[0:val]:
            Slides.text_box(slide, scores[i].name, left + spacing * i, top, spacing,
                            font_size=18, alignment=PP_ALIGN.CENTER, bold=True)
            if not isinstance(scores[i].score, str):
                fmt = (category.max_format if is_max else category.avg_format) % scores[i].score
            else:
                fmt = scores[i].score
            Slides.text_box(slide, "%s %s" % (fmt, category.unit), left + spacing * i, top + 1.3, spacing,
                            font_size=14, alignment=PP_ALIGN.CENTER)
            pic_path = 'data/pics/%s.jpg' % self.players[scores[i].name]
            if os.path.isfile(pic_path):
                slide.shapes.add_picture(pic_path, Inches(left + 0.4 + spacing * i), Inches(top + 0.45),
                                         height=Inches(0.8))

    def add_top_three_table(self, scores, slide, is_max, category, left, top):
        pos = [(0, '1st'), (1, '2nd'), (2, '3rd')]
        val = min(3, len(scores))
        for i, name in pos[0:val]:
            tx_box = slide.shapes.add_textbox(left + Inches(3*i), top + Inches(0), Inches(2.5), Inches(0.7))
            tf = tx_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER

            tx_box = slide.shapes.add_textbox(left + Inches(3*i), top + Inches(0.75), Inches(2.5), Inches(0.7))
            tf = tx_box.text_frame
            tf.text = scores[i].name
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            tx_box = slide.shapes.add_textbox(left + Inches(3*i), top + Inches(1.5), Inches(2.5), Inches(0.7))
            tf = tx_box.text_frame
            if not isinstance(scores[i].score, str):
                fmt = (category.max_format if is_max else category.avg_format) % scores[i].score
            else:
                fmt = scores[i].score
            tf.text = "%s %s" % (fmt, category.unit)
            p = tf.paragraphs[0]
            p.font.size = Pt(22)
            p.alignment = PP_ALIGN.CENTER

            pic_path = 'data/pics/%s.jpg' % self.players[scores[i].name]
            if os.path.isfile(pic_path):
                slide.shapes.add_picture(pic_path, left + Inches(3*i), top + Inches(2.5), height=Inches(2.5))

    def add_draft_suggestion(self, heroes, players, compositions):
        slide = self.add_slide(5, 152, 251, 152)
        slide.shapes.title.text = "Draft Suggestion"
        i = 0
        Slides.text_box(slide, "%.2f" % compositions[0][0], 8, 2.8, bold=True,
                        width=1.5, font_size=24, alignment=PP_ALIGN.CENTER)
        for j in range(1, min(len(compositions), 8)):
            Slides.text_box(slide, "%.2f" % compositions[j][0], 8, 4.5 + (j - 1) * 0.4, bold=True,
                            width=1.5, font_size=15, alignment=PP_ALIGN.CENTER)
        for ri, r in roles().items():
            Slides.text_box(slide, r, 0.5 + 1.5 * i, 1.5, width=1.5, font_size=14, bold=True, alignment=PP_ALIGN.CENTER)
            pic_path = 'data/heroes/%i.jpg' % compositions[0][1][ri-1][1]
            slide.shapes.add_picture(pic_path, Inches(0.7 + 1.5 * i), Inches(1.85), height=Inches(0.55))
            Slides.text_box(slide, heroes[compositions[0][1][ri-1][1]], 0.5 + 1.5 * i, 2.5,
                            width=1.5, font_size=16, alignment=PP_ALIGN.CENTER)

            pic_path = 'data/pics/%i.jpg' % compositions[0][1][ri-1][2]
            if os.path.isfile(pic_path):
                slide.shapes.add_picture(pic_path, Inches(0.85 + 1.5 * i), Inches(3), height=Inches(0.7))
            Slides.text_box(slide, players[compositions[0][1][ri-1][2]], 0.5 + 1.5 * i, 3.8,
                            width=1.5, font_size=16, alignment=PP_ALIGN.CENTER)

            for j in range(1, min(len(compositions), 8)):
                pic_path = 'data/heroes/%i.jpg' % compositions[j][1][ri - 1][1]
                slide.shapes.add_picture(pic_path, Inches(0.5 + 1.5 * i), Inches(4.5 + (j - 1) * 0.4),
                                         height=Inches(0.3))
                pic_path = 'data/pics/%i.jpg' % compositions[j][1][ri - 1][2]
                if os.path.isfile(pic_path):
                    slide.shapes.add_picture(pic_path, Inches(1.12 + 1.5 * i), Inches(4.5 + (j - 1) * 0.4),
                                             height=Inches(0.3))
            i += 1

    def add_fantasy_ranking(self, fantasy_scores):
        slide = self.add_slide(5, 104, 252, 255)
        slide.shapes.title.text = 'Current Fantasy Ranking'
        i = 0
        top = 2
        left = 0.3
        pic_size = 0.32
        column = 0.7
        row = 0.36

        Slides.text_box(slide, 'PnKasino Player', left, top - left, font_size=12, bold=True)
        Slides.text_box(slide, 'SCORE', left + 1.3, top - left, font_size=12, bold=True)
        Slides.text_box(slide, 'Reward', left + 2, top - left, font_size=12, bold=True)
        Slides.text_box(slide, 'Prize', left + 2.7, top - left, font_size=12, bold=True)
        Slides.text_box(slide, 'Team Cost', left + 3.4, top - left, font_size=12)

        Slides.text_box(slide, 'Hard Carry', left + 4.3, top - left, font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
        Slides.text_box(slide, 'Mid', left + 4.3 + 1 * (column + pic_size), top - left, font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
        Slides.text_box(slide, 'Offlane', left + 4.3 + 2 * (column + pic_size), top - left, font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
        Slides.text_box(slide, 'Support', left + 4.3 + 3 * (column + pic_size), top - left, font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
        Slides.text_box(slide, 'Hard Support', left + 4.3 + 4 * (column + pic_size), top - left, font_size=12, bold=True, alignment=PP_ALIGN.CENTER)

        for player in fantasy_scores:
            name = player['real_name'] if player['real_name'] is not None else player['name']
            Slides.text_box(slide, name, left, top + row * i, font_size=12, bold=True)
            Slides.text_box(slide, "%.2f" % player['total_score'], left + 1.3, top + row * i, font_size=12, bold=True)
            Slides.text_box(slide, "%i ₭" % player['earnings'], left + 2, top + row * i, font_size=12, bold=True)
            if player['bonus'] > 0:
                Slides.text_box(slide, "%i ₭" % player['bonus'], left + 2.7, top + row * i, font_size=12, bold=True)
            Slides.text_box(slide, "%i ₭" % player['cost'], left + 3.4, top + row * i, font_size=12)
            j = 0
            for pos in ['hard_carry', 'mid', 'offlane', 'support', 'hard_support']:
                if pos in player['team']:
                    if player['silver'] == j + 1:
                        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                               Inches(left + 4.437 + j * (column + pic_size)),
                                               Inches(top + row * i - 0.013),
                                               Inches(pic_size + 0.026),
                                               Inches(pic_size + 0.026))
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(192, 192, 192)
                        shape.line.color.rgb = RGBColor(192, 192, 192)

                    if player['gold'] == j + 1:
                        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                               Inches(left + 4.437 + j * (column + pic_size)),
                                               Inches(top + row * i - 0.013),
                                               Inches(pic_size + 0.026),
                                               Inches(pic_size + 0.026))
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(232, 170, 0)
                        shape.line.color.rgb = RGBColor(232, 170, 0)

                    pic_path = 'data/pics/%s.jpg' % self.players[player['team'][pos]['card']]
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(left + 4.45 + j * (column + pic_size)),
                                                 Inches(top + row * i), height=Inches(pic_size))

                    if player['team'][pos]['points'] > 0:
                        Slides.text_box(slide, "%.2f" % player['team'][pos]['points'],
                                        left + 4.45 + pic_size + j * (column + pic_size),
                                        top + row * i, font_size=12,
                                        bold=player['silver'] == (j + 1) or player['gold'] == (j + 1),
                                        color=RGBColor(232, 170, 0) if player['gold'] == j + 1 else
                                        (RGBColor(128, 128, 128) if player['silver'] == j + 1 else None))
                j += 1
            i += 1

    def add_fantasy_slide(self, fantasy_values, role):
        slide = self.add_slide(5, 123, 111, 255)
        slide.shapes.title.text = 'Score for the period (%s)' % role.title()
        y = 2
        x = 8
        left = 0.6
        player_size = 1
        pic_size = 0.8
        spacing = 1.82
        top = 1.5
        row_width = 2.5
        column_width = 1.1
        values = sorted([
            {
                'player': p,
                'coins': v[role] / 500,
                'silver': v['silver'][role] if role in v['silver'] else 0,
                'gold': v['gold'][role] if role in v['gold'] else 0
            } for p, v in fantasy_values.items() if v[role] > 0],
                        key=lambda e: e['coins'], reverse=True)
        for i in range(x):
            for j in range(y):
                if j * x + i < len(values):
                    c = values[j * x + i]
                    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + column_width * i),
                                           Inches(top + row_width * j),
                                           Inches(player_size),
                                           Inches(row_width - 0.1))
                    Slides.text_box(slide, c['player'], left + column_width * i, top + row_width * j + 0.2,
                                    width=player_size, font_size=12, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, "%.2f" % c['coins'], left + column_width * i,
                                    top + row_width * j + spacing - 0.2,
                                    width=player_size, font_size=16, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, "%.2f" % (c['silver'] + c['coins']), left + column_width * i,
                                    top + row_width * j + spacing + 0.1, color=RGBColor(192, 192, 192),
                                    width=player_size, font_size=12, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, "%.2f" % (c['gold'] + c['silver'] + c['coins']), left + column_width * i,
                                    top + row_width * j + spacing + 0.3, color=RGBColor(255, 215, 0),
                                    width=player_size, font_size=12, alignment=PP_ALIGN.CENTER, bold=True)
                    pic_path = 'data/pics/%s.jpg' % self.players[c['player']]
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(left + 0.1 + column_width * i),
                                                 Inches(top + row_width * j + 0.8), height=Inches(pic_size))

    def add_fantasy_data(self, fantasy_values, role):
        slide = self.add_slide(5, 255, 111, 123)
        slide.shapes.title.text = 'New Card PnKoins (%s)' % role.title()
        y = 2
        x = 9
        left = 0.1
        player_size = 1
        pic_size = 0.8
        spacing = 1.82
        top = 1.5
        row_width = 2.5
        column_width = 1.1
        values = sorted([v for v in fantasy_values if v['position'] == role],
                        key=lambda e: e['current_value'], reverse=True)
        for i in range(x):
            for j in range(y):
                if j * x + i < len(values):
                    c = values[j * x + i]
                    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + column_width * i),
                                           Inches(top + row_width * j),
                                           Inches(player_size),
                                           Inches(row_width - 0.1))
                    Slides.text_box(slide, c['name'], left + column_width * i, top + row_width * j + 0.2,
                                    width=player_size, font_size=12, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, "%i ₭" % c['current_value'], left + column_width * i,
                                    top + row_width * j + spacing - 0.23,
                                    width=player_size, font_size=16, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, "%i ₭" % c['old_value'], left + column_width * i,
                                    top + row_width * j + spacing + 0.07,
                                    width=player_size, font_size=11, alignment=PP_ALIGN.CENTER, bold=False)
                    if c['variation'] != 0:
                        Slides.text_box(slide, "%s%.1f %%" % ('+' if c['variation'] > 0 else '', c['variation'] * 100),
                                        left + column_width * i,
                                        top + row_width * j + spacing + 0.3,
                                        width=player_size, font_size=11, alignment=PP_ALIGN.CENTER, bold=True)

                    pic_path = 'data/pics/%s.jpg' % self.players[c['name']]
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(left + 0.1 + column_width * i),
                                                 Inches(top + row_width * j + 0.7), height=Inches(pic_size))

    def add_best_team(self, best_team, text, matches):
        slide = self.add_slide(5, 152, 251, 152) if text == 'Best' else self.add_slide(5, 255, 60, 60)
        slide.shapes.title.text = "%s Team by Hero Performance" % text
        i = 0
        Slides.text_box(slide, "Minimum %i or more matches for that hero on that role." % matches,
                        1, 0.1, font_size=9)
        for r in roles().values():
            if len(best_team[r]) > 0:
                pic_path = 'data/heroes/%s.jpg' % best_team[r][0]['hero_id']
                slide.shapes.add_picture(pic_path, Inches(0.5 + 1.8 * i), Inches(2), height=Inches(0.8))
                tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(3), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = best_team[r][0]['hero_name']
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].font.size = Pt(15)
                tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(3.6), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = '%.2f' % best_team[r][0]['rating']
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(1.5), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = '%s' % best_team[r][0]['role']
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                for j in range(1, min(len(best_team[r]), 5)):
                    pic_path = 'data/heroes/%s.jpg' % best_team[r][j]['hero_id']
                    slide.shapes.add_picture(pic_path, Inches(0.5 + 1.8 * i), Inches(4.5 + (j - 1) * 0.7),
                                             height=Inches(0.5))
                    tx_box = slide.shapes.add_textbox(Inches(1.6 + 1.8 * i), Inches(4.5 + (j - 1) * 0.7), Inches(1.6),
                                                      Inches(0.4))
                    tf = tx_box.text_frame
                    tf.text = '%.2f' % best_team[r][j]['rating']
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.size = Pt(16)
                    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            i += 1

    def add_rivals(self, rivals):
        slide = self.add_slide(5, 104, 222, 255)
        title_shape = slide.shapes.title
        title_shape.text = '%s Rivals' % self.team_name
        Slides.text_box(slide, "Rivals with 2 or more matches against %s." % self.team_name,
                        1, 0.1, font_size=9)

        headers = ['Rival Name', 'ID', 'Wins', 'Matches', 'Win Rate']
        keys = ['name', 'id', 'wins', 'matches', 'wr']
        formats = ['%s', '%s', '%s', '%s', '%.2f %%']
        widths = [4, 2, 1, 1, 1]
        Slides.create_table(slide, rivals, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths, hyperlink=[1], hyperlink_type='players')

    def add_trios(self, trios, text, matches):
        slide = self.add_slide(5, 255, 105, 180)
        title_shape = slide.shapes.title
        title_shape.text = '%s 15 %s Trios' % (text, self.team_name)
        Slides.text_box(slide, "Minimum %i or more matches by these trios on the same match." % matches,
                        1, 0.1, font_size=9)

        headers = ['Players', 'Rating', 'Wins', 'Matches', 'Win Rate']
        keys = ['players', 'rating', 'wins', 'matches', 'wr']
        formats = ['%s', '%.2f', '%s', '%s', '%.2f %%']
        widths = [5, 1, 1, 1, 1]
        Slides.create_table(slide, trios, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths)

    def add_lane_partners(self, partners, text, matches):
        slide = self.add_slide(5, 180, 140, 255)
        slide.shapes.title.text = "%s %s Lane Partners" % (self.team_name, text)
        Slides.text_box(slide, "Minimum %i or more matches by these players on the same lane." % matches,
                        1, 0.1, font_size=9)
        headers = ['Players', 'Rating', 'Wins', 'Losses', 'Matches', 'Win Rate']
        keys = ['lane', 'rating', 'wins', 'losses', 'matches', 'wr']
        formats = ['%s', '%.2f', '%s', '%s', '%s', '%.2f %%']
        widths = [4, 1, 1, 1, 1, 1]
        Slides.create_table(slide, partners, headers, keys, formats, Inches(0.5), Inches(1.5), Inches(9), 1, 13, 15,
                            widths=widths)

    def add_hero_lanes(self, couples, text, matches):
        slide = self.add_slide(5, 220, 125, 180)
        slide.shapes.title.text = "%s %s Heroes on Same Lane" % (text, self.team_name)
        Slides.text_box(slide, "Minimum %i or more matches by these heroes on the same lane." % matches,
                        1, 0.1, font_size=9)
        y = 5
        x = 2
        left = 0.6
        player_size = 1
        pic_size = 0.45
        spacing = 1.8
        line_space = 1.1
        top = 1.6
        column_width = 4.8
        for i in range(x):
            for j in range(y):
                if i * y + j < len(couples):
                    c = couples[i * y + j]
                    Slides.text_box(slide, c['hero_names'][0], left + column_width * i, line_space * j + top,
                                    width=player_size, font_size=11, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, c['hero_names'][1], left + player_size + spacing + column_width * i,
                                    line_space * j + top,
                                    width=player_size, font_size=11, alignment=PP_ALIGN.CENTER, bold=True)
                    pic_path = 'data/heroes/%s.jpg' % c['hero_ids'][0]
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(left + 0.1 + column_width * i),
                                                 Inches(line_space * j + top + 0.35), height=Inches(pic_size))
                    pic_path = 'data/heroes/%s.jpg' % c['hero_ids'][1]
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path,
                                                 Inches(left + 0.1 + player_size + spacing + column_width * i),
                                                 Inches(line_space * j + top + 0.35), height=Inches(pic_size))
                    Slides.text_box(slide, '%.2f' % c['rating'], left + column_width * i + player_size,
                                    line_space * j + top,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=24, bold=True)
                    Slides.text_box(slide, '%.2f %%' % c['wr'], left + column_width * i + player_size,
                                    line_space * j + top + 0.4,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=14)
                    Slides.text_box(slide, '%s matches (%s - %s)' % (c['matches'], c['wins'], c['matches'] - c['wins']),
                                    left + column_width * i + player_size, line_space * j + top + 0.7,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=12)

    def add_spammers(self, couples):
        slide = self.add_slide(5, 234, 185, 90)
        slide.shapes.title.text = "%s Hero Spammers" % self.team_name
        y = 5
        x = 2
        left = 0.6
        player_size = 1
        pic_size = 0.8
        spacing = 2
        top = 1.2
        column_width = 4.8
        for i in range(x):
            for j in range(y):
                if i * y + j < len(couples):
                    c = couples[i * y + j]
                    Slides.text_box(slide, c['player_name'], left + column_width * i, top * (1 + j),
                                    width=player_size, font_size=14, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, c['hero_name'], left + player_size + spacing + column_width * i,
                                    top * (1 + j),
                                    width=player_size, font_size=14, alignment=PP_ALIGN.CENTER, bold=True)
                    pic_path = 'data/pics/%s.jpg' % c['player_id']
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(left + 0.1 + column_width * i),
                                                 Inches(top * (1 + j) + 0.35), height=Inches(pic_size))
                    pic_path = 'data/heroes/%s.jpg' % c['hero_id']
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path,
                                                 Inches(left + 0.1 + player_size + spacing + column_width * i),
                                                 Inches(top * (1 + j) + 0.45), height=Inches(pic_size - 0.2))
                    Slides.text_box(slide, '%s' % c['matches'], left + column_width * i + player_size,
                                    top * (1 + j) + 0.1,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=24, bold=True)
                    Slides.text_box(slide, '%.2f %%' % c['wr'], left + column_width * i + player_size,
                                    top * (1 + j) + 0.55,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=16)
                    Slides.text_box(slide, 'rating %.2f, (%s - %s)' % (c['rating'], c['wins'], c['losses']),
                                    left + column_width * i + player_size, top * (1 + j) + 0.9,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=12)

    def add_hero_player_couples(self, couples, text, matches):
        slide = self.add_slide(5, 122, 105, 234)
        slide.shapes.title.text = "%s %s Hero-Player Pairs" % (self.team_name, text)
        Slides.text_box(slide, "Minimum %i or more matches by these players with these heroes." % matches,
                        1, 0.1, font_size=9)
        y = 5
        x = 2
        left = 0.6
        player_size = 1
        pic_size = 0.8
        spacing = 2
        top = 1.2
        column_width = 4.8
        for i in range(x):
            for j in range(y):
                if i * y + j < len(couples):
                    c = couples[i * y + j]
                    Slides.text_box(slide, c['player_name'], left + column_width * i, top * (1 + j),
                                    width=player_size, font_size=14, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, c['hero_name'], left + player_size + spacing + column_width * i,
                                    top * (1 + j),
                                    width=player_size, font_size=14, alignment=PP_ALIGN.CENTER, bold=True)
                    pic_path = 'data/pics/%s.jpg' % c['player_id']
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(left + 0.1 + column_width * i),
                                                 Inches(top * (1 + j) + 0.35), height=Inches(pic_size))
                    pic_path = 'data/heroes/%s.jpg' % c['hero_id']
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path,
                                                 Inches(left + 0.1 + player_size + spacing + column_width * i),
                                                 Inches(top * (1 + j) + 0.45), height=Inches(pic_size - 0.2))
                    Slides.text_box(slide, '%.2f' % c['rating'], left + column_width * i + player_size, top * (1 + j),
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=24, bold=True)
                    Slides.text_box(slide, '%.2f %%' % c['wr'], left + column_width * i + player_size,
                                    top * (1 + j) + 0.5,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=16)
                    Slides.text_box(slide, '%s matches (%s - %s)' % (c['matches'], c['wins'], c['losses']),
                                    left + column_width * i + player_size, top * (1 + j) + 0.9,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=12)

    def add_couples(self, couples, text, matches):
        inv_p = {v: k for k, v in self.players.items()}
        slide = self.add_slide(5, 255, 105, 180)
        slide.shapes.title.text = "%s %s Couples" % (self.team_name, text)
        Slides.text_box(slide, "Minimum %i or more matches by these players on the same match." % matches,
                        1, 0.1, font_size=9)
        y = 5
        x = 2
        left = 0.6
        player_size = 1
        pic_size = 0.8
        spacing = 2
        top = 1.2
        column_width = 4.8
        for i in range(x):
            for j in range(y):
                if i * y + j < len(couples):
                    c = couples[i * y + j]
                    Slides.text_box(slide, inv_p[c['p1']], left + column_width * i, top * (1 + j),
                                    width=player_size, font_size=14, alignment=PP_ALIGN.CENTER, bold=True)
                    Slides.text_box(slide, inv_p[c['p2']], left + player_size + spacing + column_width * i,
                                    top * (1 + j),
                                    width=player_size, font_size=14, alignment=PP_ALIGN.CENTER, bold=True)
                    pic_path = 'data/pics/%s.jpg' % c['p1']
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(left + 0.1 + column_width * i),
                                                 Inches(top * (1 + j) + 0.35), height=Inches(pic_size))
                    pic_path = 'data/pics/%s.jpg' % c['p2']
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path,
                                                 Inches(left + 0.1 + player_size + spacing + column_width * i),
                                                 Inches(top * (1 + j) + 0.35), height=Inches(pic_size))
                    Slides.text_box(slide, '%.2f' % c['rating'], left + column_width * i + player_size, top * (1 + j),
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=24, bold=True)
                    Slides.text_box(slide, '%.2f %%' % c['wr'], left + column_width * i + player_size,
                                    top * (1 + j) + 0.5,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=16)
                    Slides.text_box(slide, '%s matches (%s - %s)' % (c['matches'], c['wins'], c['matches'] - c['wins']),
                                    left + column_width * i + player_size, top * (1 + j) + 0.9,
                                    width=spacing, alignment=PP_ALIGN.CENTER, font_size=12)

    def add_best_team_by_player(self, best_team, text, matches):
        slide = self.add_slide(5, 152, 251, 152) if text == 'Best' else self.add_slide(5, 255, 60, 60)
        slide.shapes.title.text = "%s Player/Hero by Position" % text
        i = 0
        Slides.text_box(slide, "Minimum %i or more matches by that player for that hero on that role." % matches,
                        1, 0.1, font_size=9)
        for r in roles().values():
            if len(best_team[r]) > 0:
                tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(1.5), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = '%s' % best_team[r][0]['role']
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                pic_path = 'data/heroes/%s.jpg' % best_team[r][0]['hero_id']
                slide.shapes.add_picture(pic_path, Inches(0.5 + 1.8 * i), Inches(2), height=Inches(0.8))

                tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(3), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = best_team[r][0]['hero_name']
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].font.size = Pt(15)

                pic_path = 'data/pics/%s.jpg' % best_team[r][0]['player_id']
                if os.path.isfile(pic_path):
                    slide.shapes.add_picture(pic_path, Inches(0.8 + 1.8 * i), Inches(3.45), height=Inches(1))

                tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(4.4), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = best_team[r][0]['player_name']
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].font.size = Pt(18)

                tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(4.85), Inches(1.6), Inches(0.4))
                tf = tx_box.text_frame
                tf.text = '%.2f' % best_team[r][0]['rating']
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                for j in range(1, min(len(best_team[r]), 5)):
                    pic_path = 'data/heroes/%s.jpg' % best_team[r][j]['hero_id']
                    slide.shapes.add_picture(pic_path, Inches(0.5 + 1.8 * i), Inches(5.6 + (j - 1) * 0.45),
                                             height=Inches(0.35))
                    pic_path = 'data/pics/%s.jpg' % best_team[r][j]['player_id']
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(1.2 + 1.8 * i), Inches(5.6 + (j - 1) * 0.45),
                                                 height=Inches(0.35))
                    tx_box = slide.shapes.add_textbox(Inches(1.6 + 1.8 * i), Inches(5.6 + (j - 1) * 0.45), Inches(1.6),
                                                      Inches(0.4))
                    tf = tx_box.text_frame
                    tf.text = '%.2f' % best_team[r][j]['rating']
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.size = Pt(16)
                    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            i += 1

    def add_results_slides(self, medals, points):
        cat = Category(0, '', max_format='%s')
        slide = self.add_slide(5, 255, 255, 224)
        title_shape = slide.shapes.title
        title_shape.text = 'Top Tier Results Compilation'
        scores = [
            TierItem(medals[i][0], str(medals[i][1]), '', medals[i][1][0] * 3 + medals[i][1][1] * 2 + medals[i][1][2])
            for i in range(3)]
        self.add_top_three_table(scores, slide, True, cat, Inches(0.7), Inches(1.8))

        slide = self.add_slide(5, 255, 255, 224)
        title_shape = slide.shapes.title
        title_shape.text = 'Top Tier Results Compilation'
        headers = ['Player', 'Gold', 'Silver', 'Bronze']
        keys = ['player', 'gold', 'silver', 'bronze']
        scores = [{'player': k, 'gold': v[0], 'silver': v[1], 'bronze': v[2]} for k, v in medals]
        formats = ['%s', '%s', '%s', '%s']
        Slides.create_table_with_text_boxes(slide, [x for x in scores if x['gold'] + x['silver'] + x['bronze'] > 10],
                                            headers, keys, formats, 2, 1.5, 6, 11, 14, line_spacing=0.19)

        slide = self.add_slide(5, 255, 255, 224)
        title_shape = slide.shapes.title
        title_shape.text = 'Top Tier Weighted Points'
        scores = [TierItem(points[i][0], str(points[i][1]), '') for i in range(3)]
        self.add_top_three_table(scores, slide, True, cat, Inches(0.7), Inches(1.8))

        slide = self.add_slide(5, 255, 255, 224)
        title_shape = slide.shapes.title
        title_shape.text = 'Top Tier Weighted Points'
        headers = ['Player', 'Points']
        keys = ['player', 'points']
        scores = [{'player': k, 'points': v} for k, v in points]
        formats = ['%s', '%s']
        Slides.create_table_with_text_boxes(slide, [x for x in scores if x['points'] > 100], headers, keys, formats, 3,
                                            1.5, 3, 11, 14, line_spacing=0.19)

    def add_achievement_slide(self, achievement, result):
        inv_p = {v: k for k, v in self.players.items()}
        slide = self.add_slide(6, 0xCC, 0xCC, 0x66)
        left = top = width = height = Inches(0.4)
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        tf = txt_box.text_frame
        tf.text = achievement.name
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True

        txt_box = slide.shapes.add_textbox(left, Inches(0.9), width, height)
        tf = txt_box.text_frame
        tf.text = achievement.description
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        txt_box = slide.shapes.add_textbox(left, Inches(1.3), width, height)
        tf = txt_box.text_frame
        tf.text = 'Achievement rate: %.2f %% (%i times in %i matches)' % (
                                            result['wr'], result['wins'], result['matches'])
        tf.paragraphs[0].font.size = Pt(12)

        if achievement.img_path is not None:
            pic_path = 'data/achievements/%s' % achievement.img_path
            if os.path.isfile(pic_path):
                slide.shapes.add_picture(pic_path, Inches(7.5), Inches(0.25), height=Inches(1.5))
        if achievement.special_description:
            txt_box = slide.shapes.add_textbox(Inches(0.2), Inches(6.55), Inches(9.5), height)
            tf = txt_box.text_frame
            tf.text = 'Hero list: %s' % sequence(achievement.heroes)
            tf.word_wrap = True
            tf.paragraphs[0].font.size = Pt(10)

        if result['wins'] > 0:
            txt_box = slide.shapes.add_textbox(left + Inches(0.45), Inches(1.8), width, height)
            tf = txt_box.text_frame
            tf.text = 'Player'
            tf.paragraphs[0].font.size = Pt(12)
            txt_box = slide.shapes.add_textbox(left + Inches(1.8), Inches(1.8), width, height)
            tf = txt_box.text_frame
            tf.text = '⭐'
            tf.paragraphs[0].font.size = Pt(14)
            txt_box = slide.shapes.add_textbox(left + Inches(2.5), Inches(1.8), width, height)
            tf = txt_box.text_frame
            tf.text = 'Match List'
            tf.paragraphs[0].font.size = Pt(12)
            i = 0
            for player_id, matches in sorted([(k, v) for k, v in result['winners'].items() if len(v) > 0],
                                             key=lambda e: len(e[1]), reverse=True)[:10]:
                pic_path = 'data/pics/%s.jpg' % player_id
                if os.path.isfile(pic_path):
                    slide.shapes.add_picture(pic_path, left, Inches(2.3 + 0.42 * i), height=Inches(0.4))
                txt_box = slide.shapes.add_textbox(left + Inches(0.45), Inches(2.3 + 0.42 * i), width, height)
                tf = txt_box.text_frame
                tf.text = inv_p[player_id]
                tf.paragraphs[0].font.size = Pt(16)
                txt_box = slide.shapes.add_textbox(left + Inches(1.8), Inches(2.3 + 0.42 * i), width, height)
                tf = txt_box.text_frame
                tf.text = '%i' % len(matches)
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.bold = True
                txt_box = slide.shapes.add_textbox(left + Inches(2.5), Inches(2.35 + 0.42 * i), width, height)
                tf = txt_box.text_frame
                Slides.hyperlink_sequence(tf.paragraphs[0], ['%i' % i for i in matches], 5)
                tf.paragraphs[0].font.size = Pt(12)
                i += 1
        else:
            txt_box = slide.shapes.add_textbox(left + Inches(0.45), Inches(2.5), width, height)
            tf = txt_box.text_frame
            tf.text = 'No player earned this achievement.'
            tf.paragraphs[0].font.size = Pt(16)

    def add_popular_vote_category_slides(self, popular_vote_category):
        slide = self.add_slide(1, 221, 160, 221)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = popular_vote_category['category']
        tf = body_shape.text_frame
        tf.text = 'Opções: '
        is_images = False
        if isinstance(popular_vote_category['options'][0], list):
            ct = 0
            is_images = True
            for item in popular_vote_category['options']:
                pic_path = 'popularvote/%s' % item[0]
                if os.path.isfile(pic_path):
                    slide.shapes.add_picture(pic_path, Inches(0.5 + (ct // 2) * 1.8), Inches(2.8 + 1.8 * (ct % 2)),
                                             height=Inches(1.5))
                ct += 1
        else:
            for item in popular_vote_category['options']:
                p = tf.add_paragraph()
                p.text = item
                p.level = 1
        if is_images:
            if 'winner' in popular_vote_category and isinstance(popular_vote_category['winner'][0], list):
                tx_box = slide.shapes.add_textbox(Inches(8), Inches(1.8), Inches(2), Inches(0.5))
                tf = tx_box.text_frame
                tf.text = "Vencedores"
                for i in range(0, len(popular_vote_category['winner'])):
                    pic_path = 'popularvote/%s' % popular_vote_category['winner'][i][0]
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(8.3), Inches(2.3 + 1.55 * i), height=Inches(1))
                    tx_box = slide.shapes.add_textbox(Inches(8.3), Inches(2.3 + 1.55 * i + 1.08), Inches(2), Inches(0.4))
                    tf = tx_box.text_frame
                    tf.text = popular_vote_category['winner'][i][1]
            elif 'winner' in popular_vote_category:
                tx_box = slide.shapes.add_textbox(Inches(8), Inches(1.8), Inches(2), Inches(0.5))
                tf = tx_box.text_frame
                tf.text = "Vencedor"
                pic_path = 'popularvote/%s' % popular_vote_category['winner'][0]
                if os.path.isfile(pic_path):
                    slide.shapes.add_picture(pic_path, Inches(8), Inches(2.5), height=Inches(1.4))
                tx_box = slide.shapes.add_textbox(Inches(8), Inches(4), Inches(2), Inches(0.5))
                tf = tx_box.text_frame
                tf.text = popular_vote_category['winner'][1]
        else:
            if 'winner' in popular_vote_category and isinstance(popular_vote_category['winner'], list):
                tx_box = slide.shapes.add_textbox(Inches(8), Inches(1.8), Inches(2), Inches(0.5))
                tf = tx_box.text_frame
                tf.text = "Vencedores"
                for i in range(0, len(popular_vote_category['winner'])):
                    pic_path = 'data/pics/%s.jpg' % self.players[popular_vote_category['winner'][i]]
                    if os.path.isfile(pic_path):
                        slide.shapes.add_picture(pic_path, Inches(8.3), Inches(2.3 + 1.55 * i), height=Inches(1))
                    tx_box = slide.shapes.add_textbox(Inches(8.3), Inches(2.3 + 1.55 * i + 1.08), Inches(2), Inches(0.4))
                    tf = tx_box.text_frame
                    tf.text = popular_vote_category['winner'][i]
            elif 'winner' in popular_vote_category:
                tx_box = slide.shapes.add_textbox(Inches(8), Inches(1.8), Inches(2), Inches(0.5))
                tf = tx_box.text_frame
                tf.text = "Vencedor"
                pic_path = 'data/pics/%s.jpg' % self.players[popular_vote_category['winner']]
                if os.path.isfile(pic_path):
                    slide.shapes.add_picture(pic_path, Inches(8), Inches(2.5), height=Inches(1.4))
                tx_box = slide.shapes.add_textbox(Inches(8), Inches(4), Inches(2), Inches(0.5))
                tf = tx_box.text_frame
                tf.text = popular_vote_category['winner']

        slide = self.add_slide(5, 221, 160, 221)
        slide.shapes.title.text = popular_vote_category['category']
        chart_data = ChartData()
        if is_images:
            chart_data.categories = map(lambda e: e[1], popular_vote_category['options'])
        else:
            chart_data.categories = popular_vote_category['options']
        v = sum(popular_vote_category['votes'])
        chart_data.add_series("%i answers" % v, [x / v for x in popular_vote_category['votes']])
        x, y, cx, cy = Inches(1), Inches(1.8), Inches(8), Inches(5.5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        tx_box = slide.shapes.add_textbox(Inches(8), Inches(7), Inches(2), Inches(0.5))
        tf = tx_box.text_frame
        tf.text = "%i respostas" % v

    def add_top_five_slides(self, top):
        slide = self.add_slide(5, 221, 160, 221)
        slide.shapes.title.text = "Dream Team by Popular Vote"
        for i in range(5):
            pic_path = 'data/pics/%s.jpg' % self.players[top[i]]
            if os.path.isfile(pic_path):
                slide.shapes.add_picture(pic_path, Inches(0.5 + 1.8 * i), Inches(2), height=Inches(1.6))
            tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(4), Inches(1.6), Inches(0.4))
            tf = tx_box.text_frame
            tf.text = top[i]
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tx_box = slide.shapes.add_textbox(Inches(0.5 + 1.8 * i), Inches(5), Inches(1.6), Inches(0.4))
            tf = tx_box.text_frame
            tf.text = str(i + 1)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def save(self):
        self.presentation.save('report_%s_%s.pptx' % (self.team_name, self.dates))

    @staticmethod
    def change_slide_color(slide, r, g, b):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    @staticmethod
    def text_box(slide, text, left, top, width=None, font_size=None, alignment=None, bold=False, color=None):
        width = width if width is not None else 1
        tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
        tf = tx_box.text_frame
        tf.text = text
        tf.paragraphs[0].font.bold = bold
        if color is not None:
            tf.paragraphs[0].font.color.rgb = color
        if font_size is not None:
            tf.paragraphs[0].font.size = Pt(font_size)
        if alignment is not None:
            tf.paragraphs[0].alignment = alignment

    @staticmethod
    def histogram_win_loss(slide, data, category_label, left, top, width, height):
        chart_data = ChartData()
        chart_data.categories = [d[category_label] for d in data]
        chart_data.add_series('Wins', [d['wins'] for d in data])
        chart_data.add_series('Losses', [d['losses'] for d in data])
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                               Inches(left), Inches(top), Inches(width), Inches(height),
                                               chart_data)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    @staticmethod
    def create_table(slide, data, headers, keys, formats, left, top, width, height, font_size, header_size,
                     widths=None, hyperlink=None, hyperlink_type='matches'):
        table_shape = slide.shapes.add_table(len(data) + 1, len(headers), left, top, width, height)
        table = table_shape.table
        for i in range(len(headers)):
            table.cell(0, i).text = headers[i]
        for i in range(len(data)):
            for j in range(len(keys)):
                if hyperlink is not None and j in hyperlink:
                    run = table.cell(i + 1, j).text_frame.paragraphs[0].add_run()
                    run.text = formats[j] % data[i][keys[j]]
                    Slides.add_open_dota_link(run, formats[j] % data[i][keys[j]], hyperlink_type=hyperlink_type)
                else:
                    if isinstance(keys[j], str):
                        table.cell(i + 1, j).text = formats[j] % data[i][keys[j]]
                    else:
                        table.cell(i + 1, j).text = formats[j] % keys[j](data[i])
        Slides.set_table_font_size(table, font_size)
        for i in range(len(keys)):
            table.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(header_size)
        if widths is not None:
            for i in range(len(widths)):
                table.columns[i].width = Inches(widths[i])

    @staticmethod
    def add_open_dota_link(run, match_id, with_text=False, hyperlink_type='matches'):
        if with_text:
            run.text = match_id
        run.hyperlink.address = 'https://www.opendota.com/%s/%s' % (hyperlink_type, match_id)

    @staticmethod
    def hyperlink_sequence(paragraph, strings, maximum=0):
        if len(strings) == 1:
            Slides.add_open_dota_link(paragraph.add_run(), strings[0], with_text=True)
        elif maximum == 0 or maximum >= len(strings):
            for s in strings[:-2]:
                Slides.add_open_dota_link(paragraph.add_run(), s, with_text=True)
                r = paragraph.add_run()
                r.text = ', '
            Slides.add_open_dota_link(paragraph.add_run(), strings[-2], with_text=True)
            r = paragraph.add_run()
            r.text = ' and '
            Slides.add_open_dota_link(paragraph.add_run(), strings[-1], with_text=True)
        else:
            for s in strings[:maximum]:
                Slides.add_open_dota_link(paragraph.add_run(), s, with_text=True)
                r = paragraph.add_run()
                r.text = ', '
            r = paragraph.add_run()
            r.text = '...'

    @staticmethod
    def create_table_with_text_boxes(slide, data, headers, keys, formats, left, top, width, font_size,
                                     header_size, widths=None, line_spacing=0.25):
        if widths is None:
            widths = [width / len(headers) for _ in range(len(headers))]
        for i in range(len(headers)):
            Slides.text_box(slide, headers[i], left + sum(widths[:i]), top, widths[i], font_size=header_size, bold=True)
        for i in range(len(data)):
            for j in range(len(keys)):
                Slides.text_box(slide, formats[j] % data[i][keys[j]], left + sum(widths[:j]) + 0.05,
                                top + 0.1 + line_spacing * (i + 1), widths[j], font_size=font_size)

    @staticmethod
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    @staticmethod
    def set_table_font_size(table, size):
        for cell in Slides.iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                paragraph.space_before = Pt(1)
                paragraph.space_after = Pt(1)
                for run in paragraph.runs:
                    run.font.size = Pt(size)
